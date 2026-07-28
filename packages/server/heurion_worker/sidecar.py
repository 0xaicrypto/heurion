"""MedSci-Sidecar rendering handlers.

Implements the Execution Plane side of the MedSci-Sidecar plugin:
- DOCX generation from Jinja2-style templates (docxtpl)
- PPTX generation from templates (python-pptx)
- Medical table rendering
- Plot rendering (matplotlib / plotly)
- Upload of generated files to object storage
"""

from __future__ import annotations

import io
import json
import logging
import os
import uuid
from pathlib import Path
from typing import Any

from heurion_worker.storage import upload_output

logger = logging.getLogger("heurion-worker.sidecar")


def _template_dir() -> Path:
    """Return the directory containing bundled system templates."""
    return Path(__file__).with_name("templates")


def _make_file_id(prefix: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:16]}"


def _tenant_prefix(payload: dict[str, Any]) -> str:
    """Extract a tenant/workspace prefix from the job payload.

    Falls back to ``default`` so the object store layout is always valid.
    Accepts both snake_case (worker-native) and camelCase (Control Plane) keys.
    """
    tenant = payload.get("tenant") or {}
    if isinstance(tenant, str):
        try:
            tenant = json.loads(tenant)
        except Exception:
            tenant = {}
    workspace_id = (
        tenant.get("workspace_id")
        or tenant.get("workspaceId")
        or tenant.get("id")
        or "default"
    )
    user_id = tenant.get("user_id") or tenant.get("userId") or "anonymous"
    return f"{workspace_id}/{user_id}"


def generate_docx(payload: dict[str, Any]) -> dict:
    """Render a Word document from a template + structured data."""
    from docxtpl import DocxTemplate

    template_id = payload.get("template_id", "case_summary")
    data = payload.get("data", {})
    output_name = payload.get("output_name") or template_id

    template_path = _template_dir() / "docx" / f"{template_id}.docx"
    if not template_path.exists():
        available = [p.stem for p in (_template_dir() / "docx").glob("*.docx")]
        raise RuntimeError(f"Template '{template_id}' not found. Available: {available}")

    tpl = DocxTemplate(str(template_path))
    tpl.render(data)

    file_id = _make_file_id("docx")
    file_name = f"{output_name}.docx"
    buf = io.BytesIO()
    tpl.save(buf)

    return upload_output(
        file_obj=buf,
        file_name=file_name,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        tenant_prefix=_tenant_prefix(payload),
        file_id=file_id,
    )


def generate_pptx(payload: dict[str, Any]) -> dict:
    """Render a PowerPoint presentation from structured slide data.

    If a template file named ``{template_id}.pptx`` exists under
    ``templates/pptx/`` it is used as a starting point; otherwise a blank
    presentation is created and populated dynamically. The payload's
    ``data`` should contain ``title``, ``subtitle``, ``presenter``,
    ``date`` and a ``slides`` array of ``{title, content}`` objects.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    template_id = payload.get("template_id", "academic_presentation")
    data = payload.get("data", {})
    output_name = payload.get("output_name") or data.get("title") or template_id

    template_path = _template_dir() / "pptx" / f"{template_id}.pptx"
    if template_path.exists():
        prs = Presentation(str(template_path))
    else:
        # Fallback: build a clean 16:9 presentation from scratch.
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slides = data.get("slides") or []
    if not slides:
        # If the caller did not provide a slide list, synthesize a single
        # summary slide from whatever fields are available.
        slides = [
            {
                "title": data.get("section_title", "Summary"),
                "content": data.get("content") or data.get("findings_html") or "-",
            }
        ]

    # Ensure a title slide exists. If the template already has one, reuse it;
    # otherwise add a title slide at the front.
    title = str(data.get("title", output_name))
    subtitle = str(data.get("subtitle", ""))
    presenter = str(data.get("presenter", ""))
    date_str = str(data.get("date", data.get("generated_at", "")))

    if len(prs.slides) == 0:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
    else:
        slide = prs.slides[0]

    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        subtitle_text = "\n".join(filter(None, [subtitle, presenter, date_str]))
        slide.placeholders[1].text = subtitle_text

    # Use a Title and Content layout for the remaining slides.
    content_layout = prs.slide_layouts[1]
    for item in slides:
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(item.get("title", ""))
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            tf.text = str(item.get("content", ""))
            # Make body text readable on a projector.
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.space_after = Pt(12)

    file_id = _make_file_id("pptx")
    file_name = f"{output_name}.pptx"
    buf = io.BytesIO()
    prs.save(buf)

    return upload_output(
        file_obj=buf,
        file_name=file_name,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        tenant_prefix=_tenant_prefix(payload),
        file_id=file_id,
    )


def render_table(payload: dict[str, Any]) -> dict:
    """Render a medical/Table 1 baseline characteristics table as DOCX."""
    from docx import Document
    from docx.shared import Inches

    title = payload.get("title", "Table 1")
    headers = payload.get("headers", ["Variable", "Value"])
    rows = payload.get("rows", [])
    output_name = payload.get("output_name") or "table_1"

    doc = Document()
    doc.add_heading(title, level=1)
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"

    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = str(h)

    for row in rows:
        cells = table.add_row().cells
        for i, cell in enumerate(row):
            if i < len(cells):
                cells[i].text = str(cell)

    file_id = _make_file_id("table")
    file_name = f"{output_name}.docx"
    buf = io.BytesIO()
    doc.save(buf)

    return upload_output(
        file_obj=buf,
        file_name=file_name,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        tenant_prefix=_tenant_prefix(payload),
        file_id=file_id,
    )


def render_plot(payload: dict[str, Any]) -> dict:
    """Render a statistical plot (KM curve, bar chart, etc.) as PNG."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plot_type = payload.get("plot_type", "bar")
    title = payload.get("title", "Plot")
    x_label = payload.get("x_label", "X")
    y_label = payload.get("y_label", "Y")
    series = payload.get("series", [{"x": [1, 2, 3], "y": [1, 4, 2], "label": "Series 1"}])

    fig, ax = plt.subplots(figsize=(8, 5))
    ax.set_title(title)
    ax.set_xlabel(x_label)
    ax.set_ylabel(y_label)

    for s in series:
        x = s.get("x", [])
        y = s.get("y", [])
        label = s.get("label", "")
        if plot_type == "scatter":
            ax.scatter(x, y, label=label)
        elif plot_type == "line":
            ax.plot(x, y, label=label)
        else:
            ax.bar(x, y, label=label)

    if len(series) > 1:
        ax.legend()

    file_id = _make_file_id("plot")
    file_name = f"{payload.get('output_name', 'plot')}.png"
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    return upload_output(
        file_obj=buf,
        file_name=file_name,
        mime_type="image/png",
        tenant_prefix=_tenant_prefix(payload),
        file_id=file_id,
    )


HANDLERS = {
    "sidecar.generate_docx": generate_docx,
    "sidecar.heurion/docx.generate_docx": generate_docx,
    "sidecar.generate_pptx": generate_pptx,
    "sidecar.heurion/pptx.generate_pptx": generate_pptx,
    "sidecar.render_table": render_table,
    "sidecar.heurion/table.render_table": render_table,
    "sidecar.render_plot": render_plot,
    "sidecar.heurion/plot.render_plot": render_plot,
}


def dispatch(job_type: str, payload: dict[str, Any]) -> dict:
    handler = HANDLERS.get(job_type)
    if not handler:
        raise RuntimeError(f"Unknown sidecar job type: {job_type}")
    return handler(payload)
