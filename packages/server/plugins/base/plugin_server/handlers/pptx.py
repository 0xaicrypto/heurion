import io
import os
import uuid
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from ..storage import upload_output


def _make_file_id() -> str:
    return f"pptx_{uuid.uuid4().hex[:16]}"


def generate_pptx(payload: dict[str, Any]) -> dict:
    """Render a PowerPoint presentation from structured slide data."""
    data = payload.get("data", {})
    output_name = payload.get("output_name") or data.get("title") or "Presentation"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = str(data.get("title", output_name))
    subtitle = str(data.get("subtitle", ""))
    presenter = str(data.get("presenter", ""))
    date_str = str(data.get("date", data.get("generated_at", "")))

    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "\n".join(filter(None, [subtitle, presenter, date_str]))

    content_layout = prs.slide_layouts[1]
    slides = data.get("slides") or []
    if not slides:
        slides = [{"title": "Summary", "content": data.get("content") or data.get("findings_html") or "-"}]

    for item in slides:
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(item.get("title", ""))
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            tf.text = str(item.get("content", ""))
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.space_after = Pt(12)

    file_id = _make_file_id()
    file_name = f"{output_name}.pptx"
    buf = io.BytesIO()
    prs.save(buf)

    tenant_prefix = os.environ.get("TENANT_PREFIX", "default/anonymous")
    return upload_output(
        file_obj=buf,
        file_name=file_name,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        tenant_prefix=tenant_prefix,
        file_id=file_id,
    )
